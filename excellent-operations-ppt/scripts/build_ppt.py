#!/usr/bin/env python3
"""Render a deck JSON into a PPTX by reusing pages from assets/template.pptx."""

from __future__ import annotations

import argparse
import json
import shutil
from pathlib import Path

from pptx import Presentation


PLACEHOLDER_MARKERS = ("这里", "内容", "标题", "XX", "0.00", "00亿", "XXXX")


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) is not None and "GROUP" in str(shape.shape_type):
            try:
                yield from iter_shapes(shape.shapes)
            except Exception:
                pass


def text_shapes(slide):
    result = []
    for shape in iter_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            result.append(shape)
    return result


def set_text(shape, text: str) -> None:
    shape.text_frame.clear()
    p = shape.text_frame.paragraphs[0]
    p.text = text


def flatten_slide_text(slide: dict) -> list[str]:
    blocks = [slide.get("title", "")]
    for point in slide.get("points", [])[:4]:
        heading = point.get("heading", "")
        body = point.get("body", [])
        if isinstance(body, str):
            body = [body]
        text = "\n".join([heading, *[str(item) for item in body[:4]]]).strip()
        if text:
            blocks.append(text)
    if slide.get("summary"):
        blocks.append(str(slide["summary"]))
    return [block for block in blocks if block]


def remove_unselected_slides(prs: Presentation, keep_indices: list[int]) -> None:
    keep = set(keep_indices)
    for idx in range(len(prs.slides), 0, -1):
        if idx in keep:
            continue
        slide_id = prs.slides._sldIdLst[idx - 1]
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def render(input_json: Path, template_pptx: Path, output_pptx: Path) -> None:
    deck = json.loads(input_json.read_text(encoding="utf-8"))
    slides = deck.get("slides", [])
    if not slides:
        raise ValueError("deck JSON must contain slides")

    selected = [int(slide.get("template_slide", 13 + i)) for i, slide in enumerate(slides)]
    if len(selected) != len(set(selected)):
        raise ValueError("build_ppt.py currently requires unique template_slide values")

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(template_pptx, output_pptx)
    prs = Presentation(str(output_pptx))
    remove_unselected_slides(prs, selected)

    for slide, spec in zip(prs.slides, slides):
        replacements = flatten_slide_text(spec)
        shapes = text_shapes(slide)
        for shape, text in zip(shapes, replacements):
            set_text(shape, text)
        for shape in shapes[len(replacements):]:
            original = shape.text_frame.text.strip()
            if any(marker in original for marker in PLACEHOLDER_MARKERS):
                set_text(shape, "")

    prs.save(str(output_pptx))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("input_json", type=Path)
    parser.add_argument("output_pptx", type=Path)
    parser.add_argument(
        "--template",
        type=Path,
        default=Path(__file__).resolve().parents[1] / "assets" / "template.pptx",
    )
    args = parser.parse_args()
    render(args.input_json, args.template, args.output_pptx)


if __name__ == "__main__":
    main()
